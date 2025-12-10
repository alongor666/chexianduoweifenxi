#!/usr/bin/env python3
"""
华安保险车险周报自动生成器 v3.0 - 重构版
完全按照用户需求的报告结构设计
作者: Alongor
日期: 2025-12-09
"""

import sys
import json
import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # 非GUI后端
import io

# 导入数据转换模块
try:
    from data_transformer import DataTransformer
    DATA_TRANSFORMER_AVAILABLE = True
except ImportError:
    print("⚠ 数据转换模块未找到，仅支持标准格式数据")
    DATA_TRANSFORMER_AVAILABLE = False

# ============ 配置常量 ============
MCKINSEY_RED = RGBColor(160, 39, 36)  # #a02724
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.8)

# 状态颜色
COLOR_GREEN = RGBColor(0, 176, 80)      # #00b050 优秀
COLOR_YELLOW = RGBColor(255, 192, 0)    # #ffc000 预警
COLOR_RED = RGBColor(192, 0, 0)         # #c00000 严重
COLOR_GRAY = RGBColor(127, 127, 127)    # 中性

# ============ 1. 数据加载与验证 ============

def load_data(data_file):
    """加载数据文件(支持 xlsx/csv/json)"""
    print(f"[1/4] 加载数据: {data_file}")

    file_path = Path(data_file)
    if not file_path.exists():
        raise FileNotFoundError(f"数据文件不存在: {data_file}")

    # 根据扩展名加载
    ext = file_path.suffix.lower()
    if ext in ['.xlsx', '.xls']:
        df = pd.read_excel(data_file)
        print(f"  ✓ Excel 文件已加载 ({len(df)} 行)")
    elif ext == '.csv':
        df = pd.read_csv(data_file)
        print(f"  ✓ CSV 文件已加载 ({len(df)} 行)")
    elif ext == '.json':
        df = pd.read_json(data_file)
        print(f"  ✓ JSON 文件已加载 ({len(df)} 行)")
    else:
        raise ValueError(f"不支持的文件格式: {ext}")

    # 如果是明细数据，自动转换
    if DATA_TRANSFORMER_AVAILABLE and 'third_level_organization' in df.columns:
        print("  检测到明细数据，正在自动转换...")
        transformer = DataTransformer()
        df = transformer.transform(df)
        print(f"  ✓ 数据转换完成 ({len(df)} 行)")

    # 验证必需字段
    required_cols = ['机构', '客户类别', '签单保费', '满期赔付率', '费用率', '变动成本率']
    missing = [c for c in required_cols if c not in df.columns]
    if missing:
        raise ValueError(f"缺少必需字段: {missing}")

    print(f"✓ 数据加载成功: {len(df)} 条有效记录")
    return df

def load_config(ref_dir):
    """加载配置文件"""
    print(f"[2/4] 加载配置: {ref_dir}")

    ref_path = Path(ref_dir)
    config = {}

    # 加载阈值配置
    thresholds_file = ref_path / 'thresholds.json'
    if thresholds_file.exists():
        with open(thresholds_file, 'r', encoding='utf-8') as f:
            config['thresholds'] = json.load(f)
        print("✓ 阈值配置已加载")
    else:
        config['thresholds'] = get_default_thresholds()
        print("⚠ 使用默认阈值配置")

    # 加载保费计划
    plans_file = ref_path / 'plans.json'
    if plans_file.exists():
        with open(plans_file, 'r', encoding='utf-8') as f:
            config['plans'] = json.load(f)
        print("✓ 保费计划已加载")
    else:
        config['plans'] = None
        print("⚠ 未找到保费计划文件，保费达成率分析将跳过")

    return config

def get_default_thresholds():
    """默认阈值配置"""
    return {
        "问题机构识别阈值": {
            "年保费未达标": 95,
            "周保费未达标": 90,
            "变动成本率超标": 95,
            "满期赔付率超标": 75,
            "费用率超标": 20,
            "出险率超标": 25,
            "案均赔款超标": 6000
        },
        "四象限基准线": {
            "年计划达成率": 100,
            "变动成本率": 90,
            "满期赔付率": 70,
            "费用率": 18,
            "出险率": 20,
            "案均赔款": 6000
        },
        "状态评价阈值": {
            "变动成本率": {
                "优秀": {"max": 85, "color": "#00b050"},
                "达标": {"min": 85, "max": 95, "color": "#92d050"},
                "预警": {"min": 95, "max": 100, "color": "#ffc000"},
                "严重": {"min": 100, "color": "#c00000"}
            }
        }
    }

# ============ 2. KPI 计算 ============

def calculate_kpis(df, config):
    """计算所有KPI指标"""
    print("[3/4] 计算KPI指标")

    kpis = {}

    # 整体汇总(年累计)
    kpis['overall'] = {
        '签单保费': df['签单保费'].sum(),
        '满期保费': df.get('满期保费', df['签单保费']).sum(),
        '已报告赔款': df['已报告赔款'].sum(),
        '费用额': df.get('费用额', df['签单保费'] * df['费用率'] / 100).sum(),
        '变动成本率': weighted_average(df, '变动成本率', '签单保费'),
        '满期赔付率': weighted_average(df, '满期赔付率', '满期保费') if '满期保费' in df.columns else weighted_average(df, '满期赔付率', '签单保费'),
        '费用率': weighted_average(df, '费用率', '签单保费'),
    }

    # 计算边际贡献额
    kpis['overall']['边际贡献额'] = kpis['overall']['签单保费'] * (1 - kpis['overall']['变动成本率'] / 100)

    # 分机构汇总
    kpis['by_org'] = aggregate_by_dimension(df, '机构')

    # 分客户类别汇总
    kpis['by_customer'] = aggregate_by_dimension(df, '客户类别')

    # 计算保费达成率(如果有计划数据)
    if config['plans'] is not None:
        kpis['achievement'] = calculate_achievement_rates(df, config['plans'])
    else:
        kpis['achievement'] = None

    # 识别问题机构/客户类别
    kpis['problems'] = identify_all_problems(kpis, config['thresholds'])

    print("✓ KPI计算完成")
    return kpis

def weighted_average(df, value_col, weight_col):
    """加权平均"""
    # 如果weight_col是Series,转为列名
    if isinstance(weight_col, pd.Series):
        weight_values = weight_col
    else:
        weight_values = df[weight_col]

    return (df[value_col] * weight_values).sum() / weight_values.sum()

def aggregate_by_dimension(df, dimension):
    """按维度聚合(如果数据已经是汇总格式,直接筛选)"""

    # 检查数据是否已经是汇总格式(包含'全部'标签)
    if dimension == '机构' and '客户类别' in df.columns:
        # 筛选出 客户类别='全部' 的数据(即按机构汇总)
        result = df[df['客户类别'] == '全部'].copy()
    elif dimension == '客户类别' and '机构' in df.columns:
        # 筛选出 机构='全部' 的数据(即按客户类别汇总)
        # 如果没有这样的行,则需要重新聚合
        if (df['机构'] == '全部').any():
            result = df[df['机构'] == '全部'].copy()
        else:
            # 需要聚合
            result = _aggregate_from_detail(df, dimension)
    else:
        # 未识别格式,尝试聚合
        result = _aggregate_from_detail(df, dimension)

    # 计算占比
    total_premium = result['签单保费'].sum()
    total_claims = result['已报告赔款'].sum()

    result['保费占比'] = result['签单保费'] / total_premium * 100
    result['赔款占比'] = result['已报告赔款'] / total_claims * 100

    if '费用金额' in result.columns:
        total_expense = result['费用金额'].sum()
        result['费用占比'] = result['费用金额'] / total_expense * 100
        result['费用占比超保费占比'] = result['费用占比'] - result['保费占比']

    return result

def _aggregate_from_detail(df, dimension):
    """从明细数据聚合"""
    agg_dict = {
        '签单保费': 'sum',
        '已报告赔款': 'sum',
    }

    # 添加可选字段
    optional_fields = {
        '满期保费': 'sum',
        '费用金额': 'sum',
        '保单件数': 'sum',
        '出险件数': 'sum',
    }

    for field, func in optional_fields.items():
        if field in df.columns:
            agg_dict[field] = func

    # 执行聚合
    result = df.groupby(dimension).agg(agg_dict).reset_index()

    # 计算率值指标(加权平均)
    for dim_value in result[dimension].unique():
        dim_data = df[df[dimension] == dim_value]
        idx = result[result[dimension] == dim_value].index[0]

        result.loc[idx, '变动成本率'] = weighted_average(dim_data, '变动成本率', '签单保费')

        if '满期保费' in dim_data.columns:
            result.loc[idx, '满期赔付率'] = weighted_average(dim_data, '满期赔付率', dim_data['满期保费'])
        else:
            result.loc[idx, '满期赔付率'] = weighted_average(dim_data, '满期赔付率', '签单保费')

        result.loc[idx, '费用率'] = weighted_average(dim_data, '费用率', '签单保费')

        if '出险率' in dim_data.columns:
            result.loc[idx, '出险率'] = weighted_average(dim_data, '出险率', '签单保费')

        if '案均赔款' in dim_data.columns:
            result.loc[idx, '案均赔款'] = weighted_average(dim_data, '案均赔款', '签单保费')

    return result

def calculate_achievement_rates(df, plans):
    """计算保费达成率"""
    achievement = {'by_org': {}, 'by_customer': {}}

    # 分机构达成率
    org_premium = df.groupby('机构')['签单保费'].sum()

    if '年度保费计划' in plans:
        for org, premium in org_premium.items():
            plan = plans['年度保费计划'].get(org, 0)
            if plan > 0:
                achievement['by_org'][org] = {
                    '年累计保费': premium,
                    '年度计划': plan,
                    '年计划达成率': premium / plan * 100
                }

    # TODO: 添加周计划达成率计算

    return achievement

def identify_all_problems(kpis, thresholds):
    """识别所有问题机构和客户类别"""
    problems = {
        'org': {
            'premium_low': [],
            'cost_high': [],
            'loss_high': [],
            'expense_high': [],
            'claim_freq_high': [],
        },
        'customer': {
            'premium_low': [],
            'cost_high': [],
            'loss_high': [],
            'expense_high': [],
        }
    }

    thresh = thresholds.get('问题机构识别阈值', {})

    # 识别问题机构
    for _, row in kpis['by_org'].iterrows():
        name = row['机构']

        if row['变动成本率'] > thresh.get('变动成本率超标', 95):
            problems['org']['cost_high'].append(name)

        if row['满期赔付率'] > thresh.get('满期赔付率超标', 75):
            problems['org']['loss_high'].append(name)

        if row['费用率'] > thresh.get('费用率超标', 20):
            problems['org']['expense_high'].append(name)

    # 识别问题客户类别
    for _, row in kpis['by_customer'].iterrows():
        name = row['客户类别']

        if row['变动成本率'] > thresh.get('变动成本率超标', 95):
            problems['customer']['cost_high'].append(name)

        if row['满期赔付率'] > thresh.get('满期赔付率超标', 75):
            problems['customer']['loss_high'].append(name)

        if row['费用率'] > thresh.get('费用率超标', 20):
            problems['customer']['expense_high'].append(name)

    return problems

# ============ 3. PPT 基础组件 ============

def create_presentation():
    """创建空白PPT"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def add_blank_slide(prs):
    """添加空白幻灯片"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_top_red_line(slide)
    return slide

def add_top_red_line(slide):
    """添加顶部红线装饰"""
    shape = slide.shapes.add_shape(
        1,  # 矩形
        0, 0,
        SLIDE_WIDTH, Inches(0.015)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MCKINSEY_RED
    shape.line.fill.background()

def add_title(slide, text, font_size=24, bold=True):
    """添加标题(问题导向)"""
    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(0.4),
        SLIDE_WIDTH - 2 * MARGIN, Inches(1.0)
    )

    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    p.font.name = '微软雅黑'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(0, 0, 0)

def get_status_color(value, indicator, thresholds):
    """根据指标值返回状态颜色"""
    # 简化版本，实际需要根据thresholds配置
    if indicator == '变动成本率':
        if value < 85:
            return COLOR_GREEN
        elif value < 95:
            return RGBColor(146, 208, 80)  # 浅绿
        elif value < 100:
            return COLOR_YELLOW
        else:
            return COLOR_RED
    elif indicator == '满期赔付率':
        if value < 65:
            return COLOR_GREEN
        elif value < 75:
            return RGBColor(146, 208, 80)
        else:
            return COLOR_RED
    elif indicator == '费用率':
        if value < 16:
            return COLOR_GREEN
        elif value < 20:
            return RGBColor(146, 208, 80)
        else:
            return COLOR_RED
    else:
        return COLOR_GRAY

# ============ 4. 报告页面生成 ============

def generate_cover_page(prs, org_name, week_num, date_str):
    """生成封面"""
    slide = add_blank_slide(prs)

    # 主标题
    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(2.5),
        SLIDE_WIDTH - 2 * MARGIN, Inches(1.5)
    )

    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{org_name}车险第{week_num}周经营分析"
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.name = '微软雅黑'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED

    # 日期
    date_box = slide.shapes.add_textbox(
        MARGIN, Inches(4.5),
        SLIDE_WIDTH - 2 * MARGIN, Inches(0.5)
    )

    tf2 = date_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = date_str
    p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p2.font.name = '微软雅黑'
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_GRAY

def generate_overview_overall_page(prs, kpis, config):
    """一、经营概览——整体(年累计结果指标及状态)"""
    slide = add_blank_slide(prs)

    overall = kpis['overall']
    problems = kpis['problems']

    # 生成问题导向标题
    title_text = generate_overall_title(overall, problems, config['thresholds'])
    add_title(slide, title_text, font_size=22)

    # 左侧：率值指标(大数字展示)
    y_start = Inches(1.8)
    x_left = MARGIN

    # TODO: 添加保费时间进度达成率(需要计划数据)

    # 变动成本率
    add_kpi_display(slide, x_left, y_start, "变动成本率",
                    f"{overall['变动成本率']:.1f}%",
                    get_status_color(overall['变动成本率'], '变动成本率', config['thresholds']))

    # 满期赔付率
    add_kpi_display(slide, x_left, y_start + Inches(1.2), "满期赔付率",
                    f"{overall['满期赔付率']:.1f}%",
                    get_status_color(overall['满期赔付率'], '满期赔付率', config['thresholds']))

    # 费用率
    add_kpi_display(slide, x_left, y_start + Inches(2.4), "费用率",
                    f"{overall['费用率']:.1f}%",
                    get_status_color(overall['费用率'], '费用率', config['thresholds']))

    # 右侧：绝对值指标
    x_right = SLIDE_WIDTH / 2 + Inches(0.5)

    # 签单保费
    add_kpi_display(slide, x_right, y_start, "签单保费",
                    f"{overall['签单保费']/10000:.0f}万元", COLOR_GRAY)

    # 边际贡献额
    add_kpi_display(slide, x_right, y_start + Inches(1.2), "边际贡献额",
                    f"{overall['边际贡献额']/10000:.0f}万元", COLOR_GRAY)

    # 已报告赔款
    add_kpi_display(slide, x_right, y_start + Inches(2.4), "已报告赔款",
                    f"{overall['已报告赔款']/10000:.0f}万元", COLOR_GRAY)

    # 费用额
    add_kpi_display(slide, x_right, y_start + Inches(3.6), "费用额",
                    f"{overall['费用额']/10000:.0f}万元", COLOR_GRAY)

def add_kpi_display(slide, x, y, label, value, color):
    """添加单个KPI展示(标签+大数字)"""
    # 标签
    label_box = slide.shapes.add_textbox(x, y, Inches(3), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = '微软雅黑'
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_GRAY

    # 数值
    value_box = slide.shapes.add_textbox(x, y + Inches(0.35), Inches(3), Inches(0.7))
    tf2 = value_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = value
    p2.font.name = 'Arial'
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = color

def generate_overall_title(overall, problems, thresholds):
    """生成整体概览页的问题导向标题"""
    issues = []

    cost_rate = overall['变动成本率']
    if cost_rate < 85:
        issues.append(f"成本控制优秀({cost_rate:.1f}%)")
    elif cost_rate < 95:
        issues.append(f"成本控制良好({cost_rate:.1f}%)")
    elif cost_rate < 100:
        issues.append(f"成本率预警({cost_rate:.1f}%)")
    else:
        issues.append(f"成本率失控({cost_rate:.1f}%)")

    # 添加问题机构
    if problems['org']['cost_high']:
        org_list = '、'.join(problems['org']['cost_high'][:3])
        issues.append(f"{org_list}成本偏高")

    if not issues:
        return "本周经营状况稳健，各项指标正常"

    return "，".join(issues)

def generate_overview_org_quadrant_page(prs, kpis, config):
    """一、经营概览——分机构(年计划达成率 vs 变动成本率 四象限图)"""
    slide = add_blank_slide(prs)

    # 标题
    title_text = "概述年、周保费计划达成率与变动成本率的问题机构"
    if kpis['achievement'] is None:
        title_text = "⚠ 无保费计划数据，暂时展示变动成本率分布"

    add_title(slide, title_text, font_size=22)

    if kpis['achievement'] is None:
        # 添加提示文字
        add_text_note(slide, MARGIN, Inches(3),
                      "请提供 references/plans.json 文件以启用保费达成率分析")
        return

    # TODO: 生成四象限散点图
    # X轴: 年计划达成率, Y轴: 变动成本率
    # 基准线: X=100%, Y=90%
    add_text_note(slide, MARGIN, Inches(3), "保费达成率四象限图待实现(需plans.json)")

def generate_premium_analysis_org_page(prs, kpis, config):
    """二、保费达成分析——分机构"""
    slide = add_blank_slide(prs)

    # 标题
    title_text = generate_premium_title(kpis, config, '机构')
    add_title(slide, title_text, font_size=20)

    # 表格展示
    # 列: 机构 | 年累计保费 | 当周保费 | 年计划达成率 | 周计划达成率
    add_text_note(slide, MARGIN, Inches(3), "保费达成分析表格待实现")

def generate_premium_title(kpis, config, dimension):
    """生成保费达成分析页标题"""
    if kpis['achievement'] is None:
        return f"⚠ 无保费计划数据，无法分析{dimension}保费达成情况"

    # TODO: 根据达成率识别问题机构
    return f"指出年计划和周计划未完成的问题{dimension}"

def generate_loss_exposure_org_bubble_page(prs, kpis, config):
    """三、损失暴露分析——分机构(气泡图)"""
    slide = add_blank_slide(prs)

    title_text = generate_loss_title(kpis['by_org'], '机构')
    add_title(slide, title_text, font_size=20)

    # 生成气泡图: X轴=满期赔付率, Y轴=赔款占比, 气泡大小=赔款占比
    data = kpis['by_org'].copy()

    # 识别问题机构
    problem_orgs = []
    for _, row in data.iterrows():
        if row['满期赔付率'] > 75 and row['赔款占比'] > row['保费占比']:
            problem_orgs.append(row['机构'])

    img = create_bubble_chart(
        data=data,
        x_col='满期赔付率',
        y_col='赔款占比',
        size_col='赔款占比',
        label_col='机构',
        x_label='满期赔付率 (%)',
        y_label='赔款占比 (%)',
        title='',
        problem_items=problem_orgs
    )

    # 将图表添加到幻灯片
    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_loss_title(data, dimension):
    """生成损失暴露分析标题"""
    # 找出满期赔付率高且赔款占比超保费占比的问题项
    problems = []
    for _, row in data.iterrows():
        if row['满期赔付率'] > 75 and row['赔款占比'] > row['保费占比']:
            problems.append(row[dimension])

    if problems:
        problem_list = '、'.join(problems[:3])
        return f"{problem_list}满期赔付率偏高且赔款占比超保费占比"
    else:
        return f"各{dimension}损失暴露情况正常"

def generate_cost_analysis_org_quadrant_page(prs, kpis, config):
    """四、变动成本分析——分机构(满期赔付率 vs 费用率 四象限图)"""
    slide = add_blank_slide(prs)

    title_text = generate_cost_title(kpis['by_org'], '机构')
    add_title(slide, title_text, font_size=20)

    # 四象限图: X轴=满期赔付率, Y轴=费用率
    data = kpis['by_org'].copy()
    thresholds = config['thresholds'].get('四象限基准线', {})

    # 识别问题机构
    problem_orgs = kpis['problems']['org']['cost_high']

    img = create_quadrant_chart(
        data=data,
        x_col='满期赔付率',
        y_col='费用率',
        label_col='机构',
        x_baseline=thresholds.get('满期赔付率', 70),
        y_baseline=thresholds.get('费用率', 18),
        x_label='满期赔付率 (%)',
        y_label='费用率 (%)',
        title='',
        problem_items=problem_orgs
    )

    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_cost_title(data, dimension):
    """生成变动成本分析标题"""
    problems = []
    for _, row in data.iterrows():
        if row['变动成本率'] > 95:
            problems.append(row[dimension])

    if problems:
        problem_list = '、'.join(problems[:3])
        return f"{problem_list}变动成本率偏高，需加强成本管控"
    else:
        return f"各{dimension}变动成本控制良好"

def generate_expense_analysis_org_quadrant_page(prs, kpis, config):
    """五、费用支出分析——分机构(费用率 vs 费用占比超保费占比 四象限图)"""
    slide = add_blank_slide(prs)

    title_text = generate_expense_title(kpis['by_org'], '机构')
    add_title(slide, title_text, font_size=20)

    # 四象限图: X轴=费用率, Y轴=费用占比超保费占比
    data = kpis['by_org'].copy()
    thresholds = config['thresholds'].get('四象限基准线', {})

    # 识别问题机构
    problem_orgs = kpis['problems']['org']['expense_high']

    img = create_quadrant_chart(
        data=data,
        x_col='费用率',
        y_col='费用占比超保费占比',
        label_col='机构',
        x_baseline=thresholds.get('费用率', 18),
        y_baseline=0,  # 费用占比超保费占比的基准线是0
        x_label='费用率 (%)',
        y_label='费用占比超保费占比 (百分点)',
        title='',
        problem_items=problem_orgs
    )

    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_expense_title(data, dimension):
    """生成费用支出分析标题"""
    problems = []
    for _, row in data.iterrows():
        if row['费用率'] > 20 and row.get('费用占比超保费占比', 0) > 0:
            problems.append(row[dimension])

    if problems:
        problem_list = '、'.join(problems[:3])
        return f"{problem_list}费用率偏高且费用占比超保费占比"
    else:
        return f"各{dimension}费用支出情况正常"

def add_text_note(slide, x, y, text):
    """添加文字说明"""
    text_box = slide.shapes.add_textbox(x, y, SLIDE_WIDTH - 2*MARGIN, Inches(1))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = '微软雅黑'
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_GRAY
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# ============ 客户类别维度页面(复用机构维度逻辑) ============

def generate_premium_analysis_customer_page(prs, kpis, config):
    """二、保费达成分析——分客户类别"""
    slide = add_blank_slide(prs)
    title_text = generate_premium_title(kpis, config, '客户类别')
    add_title(slide, title_text, font_size=20)
    add_text_note(slide, MARGIN, Inches(3), "保费达成分析表格待实现(需plans.json)")

def generate_loss_exposure_customer_bubble_page(prs, kpis, config):
    """三、损失暴露分析——分客户类别(气泡图)"""
    slide = add_blank_slide(prs)
    title_text = generate_loss_title(kpis['by_customer'], '客户类别')
    add_title(slide, title_text, font_size=20)

    data = kpis['by_customer'].copy()
    problem_items = []
    for _, row in data.iterrows():
        if row['满期赔付率'] > 75 and row['赔款占比'] > row['保费占比']:
            problem_items.append(row['客户类别'])

    img = create_bubble_chart(
        data=data, x_col='满期赔付率', y_col='赔款占比', size_col='赔款占比',
        label_col='客户类别', x_label='满期赔付率 (%)', y_label='赔款占比 (%)',
        title='', problem_items=problem_items
    )
    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_loss_exposure_org_secondary_page(prs, kpis, config):
    """三、损失暴露分析——二级指标分机构(出险率 vs 案均赔款)"""
    slide = add_blank_slide(prs)

    data = kpis['by_org'].copy()
    problems = []
    for _, row in data.iterrows():
        if row.get('出险率', 0) > 25 or row.get('案均赔款', 0) > 6000:
            problems.append(row['机构'])

    if problems:
        problem_list = '、'.join(problems[:3])
        title_text = f"{problem_list}出险率或案均赔款偏高"
    else:
        title_text = "各机构出险率和案均赔款正常"

    add_title(slide, title_text, font_size=20)

    thresholds = config['thresholds'].get('四象限基准线', {})

    img = create_quadrant_chart(
        data=data, x_col='出险率', y_col='案均赔款', label_col='机构',
        x_baseline=thresholds.get('出险率', 20), y_baseline=thresholds.get('案均赔款', 6000),
        x_label='出险率 (%)', y_label='案均赔款 (元)', title='', problem_items=problems
    )
    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_loss_exposure_customer_secondary_page(prs, kpis, config):
    """三、损失暴露分析——二级指标分客户类别"""
    slide = add_blank_slide(prs)

    data = kpis['by_customer'].copy()
    problems = []
    for _, row in data.iterrows():
        if row.get('出险率', 0) > 25 or row.get('案均赔款', 0) > 6000:
            problems.append(row['客户类别'])

    if problems:
        problem_list = '、'.join(problems[:3])
        title_text = f"{problem_list}出险率或案均赔款偏高"
    else:
        title_text = "各客户类别出险率和案均赔款正常"

    add_title(slide, title_text, font_size=20)

    thresholds = config['thresholds'].get('四象限基准线', {})

    img = create_quadrant_chart(
        data=data, x_col='出险率', y_col='案均赔款', label_col='客户类别',
        x_baseline=thresholds.get('出险率', 20), y_baseline=thresholds.get('案均赔款', 6000),
        x_label='出险率 (%)', y_label='案均赔款 (元)', title='', problem_items=problems
    )
    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_cost_analysis_customer_quadrant_page(prs, kpis, config):
    """四、变动成本分析——分客户类别"""
    slide = add_blank_slide(prs)
    title_text = generate_cost_title(kpis['by_customer'], '客户类别')
    add_title(slide, title_text, font_size=20)

    data = kpis['by_customer'].copy()
    thresholds = config['thresholds'].get('四象限基准线', {})
    problem_items = kpis['problems']['customer']['cost_high']

    img = create_quadrant_chart(
        data=data, x_col='满期赔付率', y_col='费用率', label_col='客户类别',
        x_baseline=thresholds.get('满期赔付率', 70), y_baseline=thresholds.get('费用率', 18),
        x_label='满期赔付率 (%)', y_label='费用率 (%)', title='', problem_items=problem_items
    )
    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

def generate_expense_analysis_customer_quadrant_page(prs, kpis, config):
    """五、费用支出分析——分客户类别"""
    slide = add_blank_slide(prs)
    title_text = generate_expense_title(kpis['by_customer'], '客户类别')
    add_title(slide, title_text, font_size=20)

    data = kpis['by_customer'].copy()
    thresholds = config['thresholds'].get('四象限基准线', {})
    problem_items = kpis['problems']['customer']['expense_high']

    img = create_quadrant_chart(
        data=data, x_col='费用率', y_col='费用占比超保费占比', label_col='客户类别',
        x_baseline=thresholds.get('费用率', 18), y_baseline=0,
        x_label='费用率 (%)', y_label='费用占比超保费占比 (百分点)',
        title='', problem_items=problem_items
    )
    add_chart_to_slide(slide, img, MARGIN, Inches(1.8),
                      width=SLIDE_WIDTH - 2*MARGIN, height=Inches(5))

# ============ 图表生成功能 ============

def create_quadrant_chart(data, x_col, y_col, label_col, x_baseline, y_baseline,
                          x_label, y_label, title, problem_items=None):
    """创建四象限散点图(返回图片字节流)"""

    # 设置中文字体
    plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False

    fig, ax = plt.subplots(figsize=(10, 6), facecolor='white')

    # 绘制基准线
    ax.axhline(y=y_baseline, color='gray', linestyle='--', linewidth=1, alpha=0.6)
    ax.axvline(x=x_baseline, color='gray', linestyle='--', linewidth=1, alpha=0.6)

    # 绘制散点
    for _, row in data.iterrows():
        label = row[label_col]
        x_val = row[x_col]
        y_val = row[y_col]

        # 判断是否为问题项
        is_problem = problem_items and label in problem_items
        color = '#c00000' if is_problem else '#a02724'  # 问题项标红

        ax.scatter(x_val, y_val, s=200, color=color, alpha=0.7, edgecolors='white', linewidth=2)

        # 标注名称(问题项必标,其他项选择性标注)
        if is_problem or len(data) <= 12:
            ax.annotate(label, (x_val, y_val),
                       fontsize=10, ha='center', va='bottom',
                       xytext=(0, 8), textcoords='offset points')

    # 设置标签
    ax.set_xlabel(x_label, fontsize=12, color='#404040')
    ax.set_ylabel(y_label, fontsize=12, color='#404040')

    # 添加象限标签(浅灰色小字)
    x_range = ax.get_xlim()
    y_range = ax.get_ylim()

    ax.text(x_range[1]*0.95, y_range[1]*0.95, '高成本\n低达成',
            fontsize=9, color='gray', ha='right', va='top', alpha=0.5)
    ax.text(x_range[0]*1.05, y_range[1]*0.95, '双落后',
            fontsize=9, color='gray', ha='left', va='top', alpha=0.5)
    ax.text(x_range[1]*0.95, y_range[0]*1.05, '双达成',
            fontsize=9, color='gray', ha='right', va='bottom', alpha=0.5)
    ax.text(x_range[0]*1.05, y_range[0]*1.05, '低成本\n低达成',
            fontsize=9, color='gray', ha='left', va='bottom', alpha=0.5)

    # 美化图表
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, alpha=0.2, linestyle='-', linewidth=0.5)

    plt.tight_layout()

    # 保存到字节流
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()

    return img_stream

def create_bubble_chart(data, x_col, y_col, size_col, label_col,
                        x_label, y_label, title, problem_items=None):
    """创建气泡图(返回图片字节流)"""

    # 设置中文字体
    plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False

    fig, ax = plt.subplots(figsize=(10, 6), facecolor='white')

    # 计算气泡大小(归一化)
    size_values = data[size_col].values
    min_size, max_size = 200, 2000
    if size_values.max() > size_values.min():
        normalized_sizes = (size_values - size_values.min()) / (size_values.max() - size_values.min())
        bubble_sizes = min_size + normalized_sizes * (max_size - min_size)
    else:
        bubble_sizes = [min_size] * len(size_values)

    # 绘制气泡
    for idx, (i, row) in enumerate(data.iterrows()):
        label = row[label_col]
        x_val = row[x_col]
        y_val = row[y_col]

        # 判断是否为问题项
        is_problem = problem_items and label in problem_items
        color = '#c00000' if is_problem else '#a02724'

        ax.scatter(x_val, y_val, s=bubble_sizes[idx],
                  color=color, alpha=0.5, edgecolors='white', linewidth=2)

        # 标注名称(问题项必标)
        if is_problem or len(data) <= 10:
            ax.annotate(label, (x_val, y_val),
                       fontsize=10, ha='center', va='center',
                       weight='bold' if is_problem else 'normal')

    # 设置标签
    ax.set_xlabel(x_label, fontsize=12, color='#404040')
    ax.set_ylabel(y_label, fontsize=12, color='#404040')

    # 美化图表
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, alpha=0.2, linestyle='-', linewidth=0.5)

    plt.tight_layout()

    # 保存到字节流
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()

    return img_stream

def add_chart_to_slide(slide, img_stream, x, y, width, height):
    """将图表图片添加到幻灯片"""
    pic = slide.shapes.add_picture(img_stream, x, y, width=width, height=height)
    return pic

# ============ 5. 主流程 ============

def generate_report(data_file, week_num, org_name, ref_dir, output_dir="."):
    """生成完整报告"""
    print("="*60)
    print("华安保险车险周报自动生成器 v3.0")
    print("="*60)
    print()

    # 1. 加载数据
    df = load_data(data_file)

    # 2. 加载配置
    config = load_config(ref_dir)

    # 3. 计算KPI
    kpis = calculate_kpis(df, config)

    # 4. 生成PPT
    print("[4/4] 生成PPT报告")
    prs = create_presentation()

    # 日期字符串
    date_str = datetime.now().strftime("%Y年%m月%d日")

    # 生成各页(完整14-15页结构)
    generate_cover_page(prs, org_name, week_num, date_str)  # 封面
    generate_overview_overall_page(prs, kpis, config)  # 一、经营概览-整体
    generate_overview_org_quadrant_page(prs, kpis, config)  # 一、经营概览-分机构
    generate_premium_analysis_org_page(prs, kpis, config)  # 二、保费达成-分机构
    generate_premium_analysis_customer_page(prs, kpis, config)  # 二、保费达成-分客户类别
    generate_loss_exposure_org_bubble_page(prs, kpis, config)  # 三、损失暴露-分机构
    generate_loss_exposure_customer_bubble_page(prs, kpis, config)  # 三、损失暴露-分客户类别
    generate_loss_exposure_org_secondary_page(prs, kpis, config)  # 三、损失暴露-二级指标分机构
    generate_loss_exposure_customer_secondary_page(prs, kpis, config)  # 三、损失暴露-二级指标分客户类别
    generate_cost_analysis_org_quadrant_page(prs, kpis, config)  # 四、变动成本-分机构
    generate_cost_analysis_customer_quadrant_page(prs, kpis, config)  # 四、变动成本-分客户类别
    generate_expense_analysis_org_quadrant_page(prs, kpis, config)  # 五、费用支出-分机构
    generate_expense_analysis_customer_quadrant_page(prs, kpis, config)  # 五、费用支出-分客户类别

    # 保存文件
    output_file = f"{org_name}_车险周报_第{week_num}周_{datetime.now().strftime('%Y%m%d')}.pptx"
    output_path = Path(output_dir) / output_file
    prs.save(str(output_path))

    print(f"✓ 报告已生成: {output_file}")
    print()
    print("="*60)
    print("✓ 报告生成完成!")
    print("="*60)

    return str(output_path)

# ============ 命令行入口 ============

if __name__ == "__main__":
    if len(sys.argv) < 5:
        print("用法: python generate_report_v3.py <数据文件> <周次> <机构名称> <配置目录>")
        print("示例: python generate_report_v3.py data.csv 49 四川分公司 ../references")
        sys.exit(1)

    data_file = sys.argv[1]
    week_num = int(sys.argv[2])
    org_name = sys.argv[3]
    ref_dir = sys.argv[4]

    try:
        generate_report(data_file, week_num, org_name, ref_dir)
    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
