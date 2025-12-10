#!/usr/bin/env python3
"""
华安保险车险周报自动生成器 - 麦肯锡风格 v2.1
单一入口脚本，整合数据验证、KPI计算、PPT生成
新增：自动数据格式识别和转换（支持明细数据）
"""

import sys
import json
import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import re

# 导入数据转换模块
try:
    from data_transformer import DataTransformer
    DATA_TRANSFORMER_AVAILABLE = True
except ImportError:
    print("⚠ 数据转换模块未找到，仅支持标准格式数据")
    DATA_TRANSFORMER_AVAILABLE = False

# 可选依赖：DuckDB（如未安装则跳过）
try:
    import duckdb
    DUCKDB_AVAILABLE = True
except ImportError:
    DUCKDB_AVAILABLE = False
    print("⚠ DuckDB 未安装，.db 文件支持已禁用")

# ============ 配置常量 ============
MCKINSEY_RED = RGBColor(160, 39, 36)  # #a02724
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.8)

# ============ 1. 数据验证与加载 ============

def validate_and_load(data_file):
    """验证并加载数据（支持 xlsx/csv/json/duckdb + 自动格式转换）"""
    print(f"[1/4] 加载数据: {data_file}")

    required_columns = [
        '机构', '客户类别', '签单保费', '满期赔付率', '费用率',
        '变动成本率', '已报告赔款', '出险率', '案均赔款'
    ]

    try:
        file_path = Path(data_file)
        file_ext = file_path.suffix.lower()

        # 根据文件类型加载数据
        if file_ext == '.xlsx' or file_ext == '.xls':
            df = pd.read_excel(data_file)
            print(f"  ✓ Excel 文件已加载 ({len(df)} 行)")

        elif file_ext == '.csv':
            # 自动检测编码和分隔符
            df = pd.read_csv(data_file, encoding='utf-8-sig')
            print(f"  ✓ CSV 文件已加载 ({len(df)} 行)")

        elif file_ext == '.json':
            # 支持两种 JSON 格式
            with open(data_file, 'r', encoding='utf-8') as f:
                data = json.load(f)

            if isinstance(data, list):
                df = pd.DataFrame(data)  # [{...}, {...}] 格式
            elif isinstance(data, dict) and 'data' in data:
                df = pd.DataFrame(data['data'])  # {"data": [...]} 格式
            else:
                df = pd.DataFrame(data)  # 尝试直接转换

            print(f"  ✓ JSON 文件已加载 ({len(df)} 行)")

        elif file_ext == '.db' or file_ext == '.duckdb':
            if not DUCKDB_AVAILABLE:
                raise ImportError("DuckDB 未安装。请运行: pip install duckdb --break-system-packages")

            # 连接 DuckDB 并查询
            conn = duckdb.connect(data_file, read_only=True)

            # 尝试查找表名（假设主表名为 'insurance_data' 或第一个表）
            tables = conn.execute("SHOW TABLES").fetchall()
            if not tables:
                raise ValueError("DuckDB 文件中没有找到表")

            table_name = 'insurance_data' if any('insurance_data' in str(t) for t in tables) else tables[0][0]
            print(f"  ℹ 使用表: {table_name}")

            df = conn.execute(f"SELECT * FROM {table_name}").df()
            conn.close()
            print(f"  ✓ DuckDB 文件已加载 ({len(df)} 行)")

        else:
            raise ValueError(f"不支持的文件格式: {file_ext}。支持格式: .xlsx, .csv, .json, .db, .duckdb")

        # ====== 新增：自动数据格式转换 ======
        # 检查是否为标准格式
        missing = [col for col in required_columns if col not in df.columns]

        if missing and DATA_TRANSFORMER_AVAILABLE:
            print(f"\n  ℹ 检测到非标准格式数据，尝试自动转换...")
            print(f"    当前字段: {list(df.columns)[:5]}... (共{len(df.columns)}个)")

            # 使用数据转换器
            transformer = DataTransformer()
            df = transformer.transform(df)

            # 再次检查
            missing = [col for col in required_columns if col not in df.columns]
            if missing:
                print(f"  ✗ 转换后仍缺少字段: {missing}")
                raise ValueError(f"无法自动转换数据格式。缺少字段: {missing}")

        elif missing:
            print(f"  ⚠ 当前字段: {list(df.columns)}")
            print(f"  ℹ 提示：可安装数据转换模块以支持自动格式转换")
            raise ValueError(f"缺少必需字段: {missing}")
        # ====== 格式转换结束 ======

        # 数据类型转换
        numeric_cols = ['签单保费', '满期赔付率', '费用率', '变动成本率',
                       '已报告赔款', '出险率', '案均赔款']
        for col in numeric_cols:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce')

        # 删除空行
        df = df.dropna(subset=['机构', '签单保费'])

        print(f"✓ 数据加载成功: {len(df)} 条有效记录")
        return df

    except Exception as e:
        print(f"✗ 数据加载失败: {e}")
        sys.exit(1)

def load_config(config_dir):
    """加载配置文件"""
    print(f"[2/4] 加载配置: {config_dir}")
    
    thresholds_file = Path(config_dir) / "thresholds.json"
    plans_file = Path(config_dir) / "plans.json"
    
    config = {}
    
    # 加载阈值配置
    if thresholds_file.exists():
        with open(thresholds_file, 'r', encoding='utf-8') as f:
            config['thresholds'] = json.load(f)
        print("✓ 阈值配置已加载")
    else:
        print("⚠ 未找到 thresholds.json，使用默认阈值")
        config['thresholds'] = get_default_thresholds()
    
    # 加载保费计划（可选）
    if plans_file.exists():
        with open(plans_file, 'r', encoding='utf-8') as f:
            config['plans'] = json.load(f)
        print("✓ 保费计划已加载")
    else:
        print("⚠ 未找到 plans.json，保费达成率将无法计算")
        config['plans'] = None
    
    return config

def get_default_thresholds():
    """返回默认阈值配置"""
    return {
        "问题机构识别阈值": {
            "年保费未达标": 95,
            "变动成本率超标": 95,
            "满期赔付率超标": 75,
            "费用率超标": 20
        },
        "四象限基准线": {
            "保费达成率": 100,
            "变动成本率": 90,
            "满期赔付率": 70,
            "费用率": 18
        }
    }

# ============ 2. KPI 计算 ============

def calculate_kpis(df, config):
    """计算所有KPI指标"""
    print("[3/4] 计算KPI指标")
    
    kpis = {}
    
    # 整体指标
    kpis['total'] = {
        '签单保费': df['签单保费'].sum(),
        '变动成本率': df['变动成本率'].mean(),
        '满期赔付率': df['满期赔付率'].mean(),
        '费用率': df['费用率'].mean(),
        '已报告赔款': df['已报告赔款'].sum(),
        '边际贡献额': df['签单保费'].sum() * (1 - df['变动成本率'].mean() / 100)
    }
    
    # 分机构汇总
    kpis['by_org'] = df.groupby('机构').agg({
        '签单保费': 'sum',
        '变动成本率': 'mean',
        '满期赔付率': 'mean',
        '费用率': 'mean',
        '已报告赔款': 'sum',
        '出险率': 'mean',
        '案均赔款': 'mean'
    }).reset_index()
    
    # 分客户类别汇总
    kpis['by_customer'] = df.groupby('客户类别').agg({
        '签单保费': 'sum',
        '变动成本率': 'mean',
        '满期赔付率': 'mean',
        '费用率': 'mean',
        '已报告赔款': 'sum',
        '出险率': 'mean',
        '案均赔款': 'mean'
    }).reset_index()
    
    # 计算保费达成率（如果有计划数据）
    if config['plans'] is not None:
        kpis['achievement'] = calculate_achievement(df, config['plans'])
    
    # 识别问题机构
    kpis['problems'] = identify_problems(kpis, config['thresholds'])
    
    print("✓ KPI计算完成")
    return kpis

def calculate_achievement(df, plans):
    """计算保费达成率"""
    # 这里简化处理，实际需要根据累计数据计算
    achievement = {}
    # 实现逻辑待补充
    return achievement

def identify_problems(kpis, thresholds):
    """识别问题机构"""
    problems = {
        'premium_low': [],  # 保费未达标机构
        'cost_high': [],    # 成本偏高机构
        'loss_high': [],    # 赔付率偏高机构
        'expense_high': []  # 费用率偏高机构
    }
    
    thresh = thresholds.get('问题机构识别阈值', {})
    
    for _, row in kpis['by_org'].iterrows():
        org = row['机构']
        
        # 检查变动成本率
        if row['变动成本率'] > thresh.get('变动成本率超标', 95):
            problems['cost_high'].append(org)
        
        # 检查满期赔付率
        if row['满期赔付率'] > thresh.get('满期赔付率超标', 75):
            problems['loss_high'].append(org)
        
        # 检查费用率
        if row['费用率'] > thresh.get('费用率超标', 20):
            problems['expense_high'].append(org)
    
    return problems

# ============ 3. PPT 生成 ============

def create_blank_slide(prs):
    """创建空白幻灯片"""
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)
    return slide

def add_top_red_line(slide):
    """添加顶部红线"""
    shape = slide.shapes.add_shape(
        1,  # 矩形
        0, 0,
        SLIDE_WIDTH, Inches(0.015)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MCKINSEY_RED
    shape.line.fill.background()

def add_title(slide, text, font_size=24):
    """添加左对齐标题"""
    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(0.4),
        SLIDE_WIDTH - 2 * MARGIN, Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = text
    title_frame.paragraphs[0].font.size = Pt(font_size)
    title_frame.paragraphs[0].font.name = '微软雅黑'
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

def generate_ppt(kpis, config, week_number, org_name, output_path):
    """生成完整PPT"""
    print(f"[4/4] 生成PPT报告")
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 封面
    create_cover_slide(prs, org_name, week_number)
    
    # 一、经营概览（2页）
    create_overview_page1(prs, kpis, config)
    create_overview_page2(prs, kpis, config)
    
    # 二、保费进度分析（2页）
    create_premium_analysis_org(prs, kpis, config)
    create_premium_analysis_customer(prs, kpis, config)
    
    # 三、变动成本分析（2页）
    create_cost_analysis_org(prs, kpis, config)
    create_cost_analysis_customer(prs, kpis, config)
    
    # 四、损失暴露分析（4页）
    create_loss_analysis_org_bubble(prs, kpis, config)
    create_loss_analysis_customer_bubble(prs, kpis, config)
    create_loss_analysis_org_quad(prs, kpis, config)
    create_loss_analysis_customer_quad(prs, kpis, config)
    
    # 五、费用支出分析（2页）
    create_expense_analysis_org(prs, kpis, config)
    create_expense_analysis_customer(prs, kpis, config)
    
    # 保存
    prs.save(output_path)
    print(f"✓ 报告已生成: {output_path}")

def create_cover_slide(prs, org_name, week_number):
    """创建封面"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(2.5),
        SLIDE_WIDTH - 2 * MARGIN, Inches(1.5)
    )
    tf = title_box.text_frame
    tf.text = f"{org_name}车险第{week_number}周经营分析"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.name = '微软雅黑'
    tf.paragraphs[0].font.bold = True
    
    # 日期
    date_str = datetime.now().strftime("%Y年%m月%d日")
    date_box = slide.shapes.add_textbox(
        MARGIN, Inches(6.5),
        SLIDE_WIDTH - 2 * MARGIN, Inches(0.5)
    )
    tf2 = date_box.text_frame
    tf2.text = date_str
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.name = '微软雅黑'
    tf2.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

def create_overview_page1(prs, kpis, config):
    """经营概览 - 整体指标"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    
    # 生成问题导向标题
    total = kpis['total']
    problems = kpis['problems']
    
    title_parts = []
    if total['变动成本率'] < 90:
        title_parts.append(f"成本控制良好({total['变动成本率']:.1f}%)")
    elif total['变动成本率'] > 95:
        title_parts.append(f"成本偏高({total['变动成本率']:.1f}%)")
    else:
        title_parts.append(f"成本可控({total['变动成本率']:.1f}%)")
    
    if problems['cost_high']:
        title_parts.append(f"{','.join(problems['cost_high'][:2])}机构成本需关注")
    
    title = "本周" + "，".join(title_parts) if title_parts else "本周经营稳健"
    add_title(slide, title)
    
    # 左侧：核心指标（大数字）
    metrics = [
        ("签单保费", f"{total['签单保费']/10000:.0f}万", "week"),
        ("变动成本率", f"{total['变动成本率']:.1f}%", ""),
        ("满期赔付率", f"{total['满期赔付率']:.1f}%", ""),
        ("费用率", f"{total['费用率']:.1f}%", "")
    ]
    
    y_pos = Inches(1.5)
    for label, value, suffix in metrics:
        # 标签
        label_box = slide.shapes.add_textbox(
            MARGIN, y_pos,
            Inches(2), Inches(0.4)
        )
        tf = label_box.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.name = '微软雅黑'
        
        # 数值（大字）
        value_box = slide.shapes.add_textbox(
            MARGIN, y_pos + Inches(0.4),
            Inches(2.5), Inches(0.8)
        )
        tf2 = value_box.text_frame
        tf2.text = value
        tf2.paragraphs[0].font.size = Pt(48)
        tf2.paragraphs[0].font.name = 'Arial'
        tf2.paragraphs[0].font.bold = True
        tf2.paragraphs[0].font.color.rgb = MCKINSEY_RED
        
        y_pos += Inches(1.2)

def create_overview_page2(prs, kpis, config):
    """经营概览 - 分机构四象限"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    
    problems = kpis['problems']
    title = f"分机构经营状况：{','.join(problems['cost_high'][:3])}机构变动成本率偏高" if problems['cost_high'] else "各机构经营状况稳健"
    add_title(slide, title)
    
    # TODO: 添加四象限图
    # 这里需要使用 matplotlib 或直接在 PPT 中绘制散点图

def create_premium_analysis_org(prs, kpis, config):
    """保费进度分析 - 分机构"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "保费进度分析 - 分机构")
    # TODO: 实现表格展示

def create_premium_analysis_customer(prs, kpis, config):
    """保费进度分析 - 分客户类别"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "保费进度分析 - 分客户类别")
    # TODO: 实现表格展示

def create_cost_analysis_org(prs, kpis, config):
    """变动成本分析 - 分机构"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    
    problems = kpis['problems']
    title = f"{','.join(problems['cost_high'][:3])}机构变动成本率偏高" if problems['cost_high'] else "各机构成本控制良好"
    add_title(slide, title)
    # TODO: 添加四象限图

def create_cost_analysis_customer(prs, kpis, config):
    """变动成本分析 - 分客户类别"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "变动成本分析 - 分客户类别")
    # TODO: 添加四象限图

def create_loss_analysis_org_bubble(prs, kpis, config):
    """损失暴露分析 - 分机构气泡图"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    
    problems = kpis['problems']
    title = f"{','.join(problems['loss_high'][:3])}机构满期赔付率偏高" if problems['loss_high'] else "各机构赔付率正常"
    add_title(slide, title)
    # TODO: 添加气泡图

def create_loss_analysis_customer_bubble(prs, kpis, config):
    """损失暴露分析 - 分客户类别气泡图"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "损失暴露分析 - 分客户类别")
    # TODO: 添加气泡图

def create_loss_analysis_org_quad(prs, kpis, config):
    """损失暴露分析 - 分机构二级指标"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "损失暴露二级指标 - 出险率与案均赔款")
    # TODO: 添加四象限图

def create_loss_analysis_customer_quad(prs, kpis, config):
    """损失暴露分析 - 分客户类别二级指标"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "损失暴露二级指标 - 分客户类别")
    # TODO: 添加四象限图

def create_expense_analysis_org(prs, kpis, config):
    """费用支出分析 - 分机构"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    
    problems = kpis['problems']
    title = f"{','.join(problems['expense_high'][:3])}机构费用率偏高" if problems['expense_high'] else "各机构费用控制良好"
    add_title(slide, title)
    # TODO: 添加四象限图

def create_expense_analysis_customer(prs, kpis, config):
    """费用支出分析 - 分客户类别"""
    slide = create_blank_slide(prs)
    add_top_red_line(slide)
    add_title(slide, "费用支出分析 - 分客户类别")
    # TODO: 添加四象限图

# ============ 主函数 ============

def main():
    if len(sys.argv) < 4:
        print("用法: python generate_report.py <数据文件> <周次> <机构名称> [配置目录]")
        print("示例: python generate_report.py data.xlsx 49 四川分公司 ../references")
        sys.exit(1)
    
    data_file = sys.argv[1]
    week_number = sys.argv[2]
    org_name = sys.argv[3]
    config_dir = sys.argv[4] if len(sys.argv) > 4 else "../references"
    
    print(f"\n{'='*60}")
    print(f"华安保险车险周报自动生成器 - 麦肯锡风格")
    print(f"{'='*60}\n")
    
    # 1. 加载数据
    df = validate_and_load(data_file)
    
    # 2. 加载配置
    config = load_config(config_dir)
    
    # 3. 计算KPI
    kpis = calculate_kpis(df, config)
    
    # 4. 生成PPT
    output_path = f"{org_name}_车险周报_第{week_number}周_{datetime.now().strftime('%Y%m%d')}.pptx"
    generate_ppt(kpis, config, week_number, org_name, output_path)
    
    print(f"\n{'='*60}")
    print(f"✓ 报告生成完成！")
    print(f"{'='*60}\n")

if __name__ == "__main__":
    main()
