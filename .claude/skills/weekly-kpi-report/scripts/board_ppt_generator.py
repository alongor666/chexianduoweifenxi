#!/usr/bin/env python3
"""
华安保险董事会级 PPT 生成器
支持麦肯锡咨询风格模板（基于客户配色）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import json
import sys
import os
from datetime import datetime

# 加载配置文件
def load_config():
    """加载 PPT 样式配置"""
    config_path = os.path.join(os.path.dirname(__file__), '..', 'config.json')
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        return {}

CONFIG = load_config()
PPT_STYLE = CONFIG.get('PPT样式', {})

# 从配置文件读取颜色（支持十六进制）
def hex_to_rgb(hex_color):
    """将十六进制颜色转换为 RGB"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

# 麦肯锡风格配色（基于客户PPT）
MCKINSEY_RED = hex_to_rgb('#a02724')      # 客户PPT主色调
BRIGHT_RED = hex_to_rgb('#c00000')        # 预警红
DARK_GRAY = hex_to_rgb('#333333')         # 标题
MID_GRAY = hex_to_rgb('#666666')          # 正文
LIGHT_GRAY = hex_to_rgb('#E5E5E5')        # 背景
WHITE = hex_to_rgb('#FFFFFF')

# 向后兼容：支持旧的华安蓝配色
HUAAN_BLUE = hex_to_rgb(PPT_STYLE.get('华安蓝', '#0052A5'))
HUAAN_GOLD = hex_to_rgb(PPT_STYLE.get('华安金', '#FFB81C'))
GRAY_TEXT = hex_to_rgb(PPT_STYLE.get('中性灰', '#595959'))
WARNING_RED = hex_to_rgb(PPT_STYLE.get('预警红', '#FF0000'))
SAFE_GREEN = hex_to_rgb(PPT_STYLE.get('安全绿', '#00A651'))

def create_title_slide(prs, week_num):
    """创建封面页 - 麦肯锡风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 顶部红色细线（麦肯锡标志）
    top_line = slide.shapes.add_shape(
        1,  # 矩形
        Inches(0), Inches(0.3),
        Inches(13.333), Inches(0.02)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = MCKINSEY_RED
    top_line.line.fill.background()
    
    # 主标题（左对齐，极简）
    title = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(10), Inches(2)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = f"华安保险四川分支车险业务周报"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    
    # 副标题
    subtitle = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),
        Inches(10), Inches(1)
    )
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = f"第 {week_num} 周经营分析"
    p.font.size = Pt(24)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 底部信息（左下角）
    footer = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),
        Inches(8), Inches(0.5)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"华安保险四川分公司车险部 | {datetime.now().strftime('%Y年%m月')}"
    p.font.size = Pt(12)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"

def create_executive_summary(prs, kpis):
    """创建经营概览页 - 麦肯锡风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部红色细线
    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0.6), Inches(13.333), Inches(0.015))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = MCKINSEY_RED
    top_line.line.fill.background()
    
    # 结论性标题（麦肯锡So What原则）
    combined_ratio = kpis['盈利能力']['综合成本率']
    premium = kpis['业务规模']['总保费_万元']
    
    title_text = f"本周保费{premium/10000:.1f}亿元，综合成本率{combined_ratio:.1f}%{'达标' if combined_ratio < 95 else '需关注'}"
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 核心数据（左侧，要点化）
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "核心数据"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    data_points = [
        f"总保费：{premium:,.0f} 万元",
        f"保单数：{kpis['业务规模']['保单数量']:,} 单",
        f"综合成本率：{combined_ratio:.1f}%",
        f"新能源车占比：{kpis['新能源车分析'].get('新能源车保费占比', 0):.1f}%"
    ]
    
    for point in data_points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)
        p.level = 0
    
    # 关键建议（右侧）
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "关键建议"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    for action in kpis['行动建议']:
        p = tf.add_paragraph()
        p.text = f"• {action.replace('⚠️', '').replace('🔋', '').replace('🚨', '').replace('📉', '').replace('✅', '').strip()}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)
        p.level = 0
    
    # 页码
    page_num = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.8), Inches(0.3))
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "2"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

def create_profitability_slide(prs, kpis):
    """创建盈利能力分析页 - 麦肯锡风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部红色细线
    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0.6), Inches(13.333), Inches(0.015))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = MCKINSEY_RED
    top_line.line.fill.background()
    
    loss_ratio = kpis['盈利能力']['平均赔付率']
    expense_ratio = kpis['盈利能力']['平均费用率']
    combined_ratio = kpis['盈利能力']['综合成本率']
    
    # 结论性标题
    if combined_ratio < 95:
        title_text = f"盈利能力稳健，综合成本率{combined_ratio:.1f}%低于行业基准"
    elif combined_ratio < 100:
        title_text = f"成本率{combined_ratio:.1f}%接近盈亏线，需加强管控"
    else:
        title_text = f"盈利承压，综合成本率{combined_ratio:.1f}%超出盈亏平衡线"
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY if combined_ratio < 100 else BRIGHT_RED
    p.font.name = "Microsoft YaHei"
    
    # 左侧：成本率拆解
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5))
    tf = left_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "成本率拆解"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(15)
    
    # 大数字展示
    p = tf.add_paragraph()
    p.text = f"{combined_ratio:.1f}%"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED if combined_ratio < 100 else BRIGHT_RED
    p.font.name = "Arial"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "综合成本率"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(20)
    
    # 拆解数据
    details = [
        f"赔付率：{loss_ratio:.1f}%",
        f"费用率：{expense_ratio:.1f}%"
    ]
    for detail in details:
        p = tf.add_paragraph()
        p.text = f"• {detail}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(5)
    
    # 右侧：高成本业务
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5))
    tf = right_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "高成本业务风险"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    high_cost = kpis['盈利能力']['高成本业务TOP3']
    if high_cost:
        for biz, count in list(high_cost.items())[:3]:
            p = tf.add_paragraph()
            p.text = f"• {biz}：{count} 单"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Microsoft YaHei"
            p.space_before = Pt(8)
    
    # 底部结论
    conclusion = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.5), Inches(1))
    tf = conclusion.text_frame
    p = tf.paragraphs[0]
    if combined_ratio < 95:
        p.text = "成本控制有效，保持现有风控策略"
    elif combined_ratio < 100:
        p.text = "建议：收紧高成本业务承保，优化费用结构"
    else:
        p.text = "建议：立即暂停高风险业务，全面审视费用率和赔付率"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 页码
    page_num = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.8), Inches(0.3))
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "3"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

def create_nev_slide(prs, kpis):
    """创建新能源车专题页 - 麦肯锡风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部红色细线
    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0.6), Inches(13.333), Inches(0.015))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = MCKINSEY_RED
    top_line.line.fill.background()
    
    nev_data = kpis['新能源车分析']
    if nev_data.get('新能源车数据') == '无':
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "新能源车业务：暂无数据"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        return
    
    nev_ratio = nev_data['新能源车保费占比']
    nev_loss = nev_data['新能源车平均赔付率']
    traditional_loss = nev_data['传统车平均赔付率']
    gap = nev_loss - traditional_loss
    
    # 结论性标题
    if gap > 20:
        title_text = f"新能源车盈利性堪忧：赔付率{nev_loss:.1f}%，高出传统车{gap:.1f}pp"
    elif gap > 10:
        title_text = f"新能源车成本偏高：赔付率{nev_loss:.1f}%，需优化定价模型"
    else:
        title_text = f"新能源车业务稳健：占比{nev_ratio:.1f}%，赔付率{nev_loss:.1f}%"
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_RED if gap > 20 else DARK_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 左侧：核心指标
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5))
    tf = left_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "新能源车核心指标"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(15)
    
    # 大数字：赔付率
    p = tf.add_paragraph()
    p.text = f"{nev_loss:.1f}%"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_RED if nev_loss > 100 else MCKINSEY_RED
    p.font.name = "Arial"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "NEV赔付率"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(20)
    
    # 其他指标
    metrics = [
        f"保费占比：{nev_ratio:.1f}%",
        f"保单数：{nev_data['新能源车保单数']:,} 单",
        f"单均保费：{nev_data['新能源车单均保费']:,.0f} 元"
    ]
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(5)
    
    # 右侧：对比分析
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5))
    tf = right_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "与传统车对比"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    comparisons = [
        f"传统车赔付率：{traditional_loss:.1f}%",
        f"赔付率差距：+{gap:.1f}pp",
        f"差距幅度：{(gap/traditional_loss*100):.1f}%"
    ]
    for comp in comparisons:
        p = tf.add_paragraph()
        p.text = f"• {comp}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)
    
    # 底部建议
    conclusion = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.5), Inches(1))
    tf = conclusion.text_frame
    p = tf.paragraphs[0]
    if gap > 20:
        p.text = "建议：立即优化NEV定价模型，考虑差异化费率；加强新能源车风险评估"
    elif gap > 10:
        p.text = "建议：调整新能源车保费基准，强化维修成本管控"
    else:
        p.text = "NEV业务风险可控，继续关注市场动态"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 页码
    page_num = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.8), Inches(0.3))
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "4"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

def create_risk_slide(prs, kpis):
    """创建风险管理页 - 麦肯锡风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部红色细线
    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0.6), Inches(13.333), Inches(0.015))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = MCKINSEY_RED
    top_line.line.fill.background()
    
    risk_data = kpis['风险指标']
    freq = risk_data['平均出险频度']
    claim_amt = risk_data['案均赔款_元']
    
    # 结论性标题
    if freq > 25:
        title_text = f"风险暴露偏高：出险频度{freq:.1f}%，需加强承保筛查"
    elif freq > 20:
        title_text = f"风险指标需关注：出险频度{freq:.1f}%，案均赔款{claim_amt:,.0f}元"
    else:
        title_text = f"风险管控有效：出险频度{freq:.1f}%处于合理区间"
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_RED if freq > 25 else DARK_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 左侧：核心风险指标
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5))
    tf = left_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "核心风险指标"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(15)
    
    # 大数字：出险频度
    p = tf.add_paragraph()
    p.text = f"{freq:.1f}%"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_RED if freq > 25 else MCKINSEY_RED
    p.font.name = "Arial"
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "平均出险频度"
    p.font.size = Pt(14)
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(20)
    
    # 其他风险数据
    metrics = [
        f"总案件数：{risk_data['总案件数']:,} 件",
        f"案均赔款：{claim_amt:,.0f} 元",
        f"高频出险占比：{risk_data['高频出险业务占比']:.1f}%"
    ]
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(5)
    
    # 右侧：高风险业务
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5))
    tf = right_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "高风险业务分布"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    high_risk = risk_data.get('高风险业务类型', {})
    if high_risk:
        for biz, count in list(high_risk.items())[:3]:
            p = tf.add_paragraph()
            p.text = f"• {biz}：{count} 件"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Microsoft YaHei"
            p.space_before = Pt(8)
    
    # 底部建议
    conclusion = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.5), Inches(1))
    tf = conclusion.text_frame
    p = tf.paragraphs[0]
    if freq > 25:
        p.text = "建议：立即强化核保审核，重点筛查高频出险业务类型"
    elif freq > 20:
        p.text = "建议：加强风险监控，定期复盘高风险业务承保决策"
    else:
        p.text = "风险控制有效，维持现有核保标准"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = MID_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 页码
    page_num = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.8), Inches(0.3))
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "5"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

def create_comparison_slide(prs, comparison_data):
    """创建周环比对比页 - 麦肯锡风格（可选）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部红色细线
    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0.6), Inches(13.333), Inches(0.015))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = MCKINSEY_RED
    top_line.line.fill.background()
    
    slide_data = comparison_data.get('幻灯片数据', {})
    title_text = slide_data.get('幻灯片标题', '周环比分析')
    
    # 提取关键变化，生成结论性标题
    evaluations = comparison_data.get('详细变化', {}).get('综合评估', [])
    if evaluations:
        title_text = evaluations[0].replace('📈', '').replace('📉', '').strip()
    
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    
    # 左侧：关键变化
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4))
    tf = left_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "关键变化"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    for change in slide_data.get('关键变化', []):
        p = tf.add_paragraph()
        p.text = f"• {change}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)
        p.level = 0
    
    # 右侧：综合评估
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4))
    tf = right_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "综合评估"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MCKINSEY_RED
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)
    
    for evaluation in slide_data.get('综合评估', []):
        clean_eval = evaluation.replace('📈', '').replace('📉', '').replace('📊', '').replace('⚠️', '').replace('✅', '').strip()
        p = tf.add_paragraph()
        p.text = f"• {clean_eval}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)
        p.level = 0
    
    # 页码
    page_num = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.8), Inches(0.3))
    tf = page_num.text_frame
    p = tf.paragraphs[0]
    p.text = "6"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

def generate_board_ppt(kpi_file, week_num, comparison_data=None):
    """
    生成完整的董事会 PPT - 麦肯锡咨询风格
    
    参数:
        kpi_file: KPI 数据文件路径
        week_num: 周次
        comparison_data: 可选的环比对比数据
    """
    with open(kpi_file, 'r', encoding='utf-8') as f:
        kpis = json.load(f)
    
    # 优先加载麦肯锡风格模板
    mckinsey_template = os.path.join(os.path.dirname(__file__), '..', 'resources', 'mckinsey_board_template.pptx')
    huaan_template = os.path.join(os.path.dirname(__file__), '..', 'resources', 'huaan_board_template.pptx')
    
    if os.path.exists(mckinsey_template):
        try:
            prs = Presentation(mckinsey_template)
            print(f"✅ 使用麦肯锡风格模板: {mckinsey_template}", file=sys.stderr)
        except:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
    elif os.path.exists(huaan_template):
        try:
            prs = Presentation(huaan_template)
            print(f"✅ 使用华安保险模板: {huaan_template}", file=sys.stderr)
        except:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
    else:
        # 创建16:9空白演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        print("⚠️ 未找到模板，使用空白16:9格式", file=sys.stderr)
    
    # 生成各页幻灯片（麦肯锡风格）
    create_title_slide(prs, week_num)
    create_executive_summary(prs, kpis)
    create_profitability_slide(prs, kpis)
    create_nev_slide(prs, kpis)
    create_risk_slide(prs, kpis)
    
    # 如果有环比数据，添加环比对比页
    if comparison_data:
        create_comparison_slide(prs, comparison_data)
    
    # 保存
    output_path = f'/mnt/user-data/outputs/华安车险周报_第{week_num}周_麦肯锡版.pptx'
    prs.save(output_path)
    
    return output_path

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "请提供 KPI 文件路径"}, ensure_ascii=False))
        sys.exit(1)
    
    kpi_file = sys.argv[1]
    week_num = sys.argv[2] if len(sys.argv) > 2 else "未知"
    comparison_file = sys.argv[3] if len(sys.argv) > 3 else None
    
    try:
        # 加载环比数据（如果提供）
        comparison_data = None
        if comparison_file and os.path.exists(comparison_file):
            with open(comparison_file, 'r', encoding='utf-8') as f:
                comparison_data = json.load(f)
            print(f"✅ 加载环比对比数据: {comparison_file}", file=sys.stderr)
        
        output_path = generate_board_ppt(kpi_file, week_num, comparison_data)
        
        result = {
            "status": "success",
            "ppt_file": output_path,
            "message": f"第 {week_num} 周董事会 PPT 生成成功"
        }
        
        if comparison_data:
            result["包含环比分析"] = True
        
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        print(json.dumps({"error": str(e)}, ensure_ascii=False))
        sys.exit(1)

if __name__ == '__main__':
    main()
